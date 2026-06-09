import streamlit as st
import sqlite3
import pandas as pd
from datetime import date, timedelta
from pptx import Presentation
import os

# ================= CONFIG =================
ORGANIZATION = "Commercial Bank of Ethiopia"
DISTRICT = "Jigjiga District IS Support"
DEFAULT_PASSWORD = "cbe@123"

STAFF_LIST = [
    "Esayase Belay", "Amanuel Desta", "Abel Alemu", "Tensay Lemessa",
    "Muna Mohammed", "Mulugeta Zewude", "Getalem Abera", "Zemen Wondesen", "Ayder Sayb"
]

LIMITED_USERS = ["Ayder Sayb", "Getalem Abera", "Zemen Wondesen", "Muna Mohammed"]

BRANCH_LIST = [
    "Jigjiga District", "Abadir Branch", "Aboker Branch", "Al-Haramain CBE-Noor Branch",
    "Al-Huda CBE Noor Branch", "Arafa CBE Noor Branch", "Ararso Branch", "Babile Branch",
    "Bonbas Branch", "Boqolemayo Branch", "Bulale Branch", "Cherety Branch", "Chinagesen Branch",
    "Degehabour Branch", "Deghale Branch", "Denan CBE Noor Branch", "Dewelle Branch", "Dollo Ado Branch",
    "Erer Branch", "Fafen Branch", "Fedis Branch", "Fiq Branch", "Fugnan Birra Branch", "Galmashira Branch",
    "Galoo CBE Noor Branch", "God-Cusbo Branch", "Gode Branch", "Hadigalla Branch", "Hakim Branch", "Harar Branch",
    "Hargelle Branch", "Hartishek Branch", "Hilal CBE Noor Branch", "Iftiin Branch", "Imam Ahmed CBE Noor Branch",
    "Jarso Branch", "Jegol Branch", "Jehedin Branch", "Jerer Branch", "Jijiga Branch", "Jinella Branch", "Kali Branch",
    "Kalub Branch", "Kebri Beya Branch", "Kebri Dehar Branch", "Kelafo Branch", "Korahe Branch", "Melka Rafu Branch",
    "Midega Branch", "Mubarek CBE Noor Branch", "Safwa CBE Noor Branch", "Sebati CBE Noor Branch", "Shebele Branch",
    "Shenkor Branch", "Shinille Branch", "Sulul Branch", "Togochale Branch", "Warder Branch", "Wilwal Branch"
]

DISTRICT_UNITS = [
    "Jigjiga District IS Support","Jijiga District HR", "Jigjiga District Directors office", "Jigjiga District Facilities",
    "Jigjiga District Digital support", "Jigjiga District Security",
    "Jigjiga District Digital Banking", "Jigjiga District Internal control",
    "Jigjiga District Internal Audit", "Jigjiga District Legal support",
    "Jigjiga District CBE-Noor IFB", "Jigjiga District Whole sale",
    "Jigjiga District Retail", "Jigjiga District Credit"
]

TASK_LIST = [
    "Network troubleshooting", "Empower ID","Western Union", "World remit",
    "Dhabshill user/ setup","Moneygram User/satup", "T24 Modification ",
    "T24 User reset", "T24 Authorization", "PC / hardware maintenance", "National ID support", "System / application support",
    "Printer / Share or Install", "User account / AD reset", "Backup / data support",
    "Install GX, Java, Ticket Writer", "Other"
]

DB_FILE = "weekly_report.db"
first_time_db = not os.path.exists(DB_FILE)

# ================= DATABASE =================
conn = sqlite3.connect(DB_FILE, check_same_thread=False)
cursor = conn.cursor()

# Create table if not exists
cursor.execute("""
CREATE TABLE IF NOT EXISTS jobs (
    id INTEGER PRIMARY KEY AUTOINCREMENT,
    job_date TEXT,
    staff_name TEXT,
    branch TEXT,
    district_unit TEXT,
    category TEXT,
    task TEXT,
    reset_for TEXT,
    description TEXT,
    status TEXT
)
""")

conn.commit()

# Ensure reset_for column exists (for old database compatibility)
try:
    cursor.execute("ALTER TABLE jobs ADD COLUMN reset_for TEXT")
    conn.commit()
except sqlite3.OperationalError:
    pass

st.markdown("""
<style>
.footer {
    position: fixed;
    left: 0; bottom: 0;
    width: 100%;
    background-color: #0E1117;
    color: grey;
    text-align: center;
    padding: 5px;
    font-size: 20px;
}
</style>
<div class="footer"> © 2026 Developed by JDIS </div> 
""", unsafe_allow_html=True)


# ================= SESSION =================
if "logged_in" not in st.session_state:
    st.session_state.logged_in = False
    st.session_state.role = None
    st.session_state.username = None
    st.session_state.page = "Dashboard"

# ================= HELPERS =================
def get_staff_name(username):
    mapping = {s.replace(" ", "").lower(): s for s in STAFF_LIST}
    return mapping.get(username, username)


def logout_button(location=""):
        if st.button("🚪 Logout", key=f"logout_{location}"):
            # Clear session safely
            st.session_state.clear()

            # Rerun app
            st.rerun()

# ================= LOGIN PAGE =================
def login_page():
    st.set_page_config(page_title="CBE JDIS Weekly Report", layout="centered")

    st.markdown("""
        <style>
        .login-box {
            max-width:400px;
            margin:auto;
            padding:40px;
            background:white;
            border-radius:12px;
            box-shadow:0 8px 20px rgba(0,0,0,0.2);
            animation: fadeIn 1.2s ease-in-out;
        }
        @keyframes fadeIn {
            from {opacity:0; transform: translateY(-20px);}
            to {opacity:1; transform: translateY(0);}
        }
        </style>
    """, unsafe_allow_html=True)

    st.markdown('<div class="login-box">', unsafe_allow_html=True)
    st.markdown("### 🔐 JDIS Login")

    with st.form("login_form"):
        username_input = st.text_input("Username").lower().replace(" ", "")
        password_input = st.text_input("Password", type="password")
        login_submit = st.form_submit_button("Login")

        if login_submit:
            user = cursor.execute(
                "SELECT * FROM users WHERE username=? AND password=?",
                (username_input, password_input)
            ).fetchone()

            if user:
                st.session_state.logged_in = True
                st.session_state.role = user[2]
                st.session_state.username = user[0]
                st.session_state.page = "🏠 Dashboard"
                st.rerun()
            else:
                st.error("Invalid credentials")

    st.markdown("</div>", unsafe_allow_html=True)


# ================= DASHBOARD =================
def dashboard_page():

    username = get_staff_name(st.session_state.username)
    role = st.session_state.role

    if st.button("🚪 Logout"):
        st.session_state.clear()
        st.rerun()

    st.markdown("""
    <style>
    .welcome {
        font-size:34px;
        font-weight:bold;
        text-align:center;
        animation: slideDown 1s ease-out;
    }
    @keyframes slideDown {
        from {opacity:0; transform: translateY(-30px);}
        to {opacity:1; transform: translateY(0);}
    }
    .kpi {
        padding:25px;
        border-radius:15px;
        text-align:center;
        color:white;
        font-weight:bold;
        transition:0.3s;
    }
    .kpi:hover {
        transform: scale(1.05);
    }
    </style>
    """, unsafe_allow_html=True)

    st.markdown(f'<div class="welcome">👋 Welcome, {username}</div>', unsafe_allow_html=True)

    df_jobs = pd.read_sql("SELECT * FROM jobs", conn)

    if role == "staff":
        df_jobs = df_jobs[df_jobs["staff_name"].apply(
            lambda x: x.replace(" ", "").lower() == st.session_state.username
        )]

    completed = len(df_jobs[df_jobs["status"]=="Completed"])
    progress = len(df_jobs[df_jobs["status"]=="In Progress"])
    total = len(df_jobs)

    col1, col2, col3 = st.columns(3)

    col1.markdown(f'<div class="kpi" style="background:green">Completed<br><h2>{completed}</h2></div>', unsafe_allow_html=True)
    col2.markdown(f'<div class="kpi" style="background:white;color:black">In Progress<br><h2>{progress}</h2></div>', unsafe_allow_html=True)
    col3.markdown(f'<div class="kpi" style="background:#1f77b4">Total Jobs<br><h2>{total}</h2></div>', unsafe_allow_html=True)

    st.subheader("📊 Task Overview")
    if not df_jobs.empty:
        st.bar_chart(df_jobs["task"].value_counts())

    st.subheader("📊 Category Overview")
    if not df_jobs.empty:
        st.bar_chart(df_jobs["category"].value_counts())

    with st.expander("🔽 View All Jobs"):
        if df_jobs.empty:
            st.info("No jobs found.")
        else:
            st.dataframe(df_jobs)


# ================= JOB ENTRY =================
def job_entry_page():

    st.header("📝 Submit Job")

    if st.button("🚪 Logout"):
        st.session_state.clear()
        st.rerun()

    st.markdown("""
    <style>
    .form-animate {
        animation: fadeInForm 0.8s ease-in-out;
    }
    @keyframes fadeInForm {
        from {opacity:0; transform: translateY(20px);}
        to {opacity:1; transform: translateY(0);}
    }
    </style>
    """, unsafe_allow_html=True)

    st.markdown('<div class="form-animate">', unsafe_allow_html=True)

    with st.form("add_form"):

        job_date = st.date_input("Date", value=date.today())

        staff_name = get_staff_name(st.session_state.username) \
            if st.session_state.role=="staff" else st.selectbox("Staff Name", STAFF_LIST)

        branch = st.selectbox("Branch", BRANCH_LIST)

        district_unit = st.selectbox("District Unit", DISTRICT_UNITS) \
            if branch=="Jigjiga District" else ""

        category_options = ["Support Remotely","In Person","Other"]
        category = st.selectbox("Job Category", category_options)

        task = st.selectbox("Task", TASK_LIST)
        reset_for = ""

        if task == "User account / AD reset":
            reset_for = st.text_input("Reset AD For (User Full Name) *")

        description_label = "Description *" if category=="Other" else "Description"
        description = st.text_area(description_label)

        status = st.selectbox("Status", ["Completed","In Progress"])

        submit = st.form_submit_button("Submit Task")

        if "job_submitted" in st.session_state:
            st.success("Job added successfully ✅")
            del st.session_state.job_submitted

        if submit:

            if category == "Other" and not description.strip():
                st.error("Description is required when category is 'Other'")
            elif task == "Other" and not description.strip():
                st.error("Description is required when task is 'Other'")
            elif task == "User account / AD reset" and not reset_for.strip():
                st.error("You must specify whose AD was reset.")
            else:
                cursor.execute("""
                   INSERT INTO jobs 
                    (job_date, staff_name, branch, district_unit, category, task, reset_for, description, status)
                    VALUES (?,?,?,?,?,?,?,?,?)
                """, (job_date, staff_name, branch, district_unit, category, task, reset_for, description, status))

                conn.commit()

                st.session_state.job_submitted = True
                st.rerun()

def update_delete_page():
    st.markdown("### 📝 Update or Delete a Job")

    df_jobs = pd.read_sql("SELECT * FROM jobs", conn)

    def normalize_username(name):
        return str(name).replace(" ", "").lower()

    if st.session_state.role == "staff":
        df_jobs = df_jobs[
            df_jobs["staff_name"].apply(
                lambda x: normalize_username(x) == st.session_state.username
            )
        ]

    if df_jobs.empty:
        st.info("No jobs found.")
        return

    job_ids = df_jobs["id"].tolist()
    selected_job_id = st.selectbox("Select Job to Edit/Delete", job_ids)

    job_row = df_jobs[df_jobs["id"] == selected_job_id].iloc[0]

    STAFF_LIST = pd.read_sql("SELECT DISTINCT staff_name FROM jobs", conn)["staff_name"].dropna().tolist()
    BRANCH_LIST = pd.read_sql("SELECT DISTINCT branch FROM jobs", conn)["branch"].dropna().tolist()
    DISTRICT_UNITS = pd.read_sql("SELECT DISTINCT district_unit FROM jobs", conn)["district_unit"].dropna().tolist()
    TASK_LIST = pd.read_sql("SELECT DISTINCT task FROM jobs", conn)["task"].dropna().tolist()

    status_options = ["In Progress", "Completed"]

    def safe_index(lst, value):
        try:
            return lst.index(value)
        except ValueError:
            return 0

    with st.form(f"job_form_{selected_job_id}"):

        job_date = st.date_input(
            "Job Date",
            pd.to_datetime(job_row["job_date"])
        )

        staff_name = st.selectbox(
            "Staff Name",
            STAFF_LIST,
            index=safe_index(STAFF_LIST, job_row["staff_name"])
        )

        branch = st.selectbox(
            "Branch",
            BRANCH_LIST,
            index=safe_index(BRANCH_LIST, job_row["branch"])
        )

        district_unit = st.selectbox(
            "District Unit",
            DISTRICT_UNITS,
            index=safe_index(DISTRICT_UNITS, job_row["district_unit"])
        )

        task = st.selectbox(
            "Task",
            TASK_LIST,
            index=safe_index(TASK_LIST, job_row["task"])
        )

        # Handle AD Reset field
        reset_for = job_row["reset_for"]

        if task == "User account / AD reset":
            reset_for = st.text_input(
                "Reset AD For (User Full Name)",
                value=job_row["reset_for"] if job_row["reset_for"] else ""
            )

        description = st.text_area(
            "Description",
            job_row["description"]
        )

        status = st.selectbox(
            "Status",
            status_options,
            index=safe_index(status_options, job_row["status"])
        )

        # ✅ Correct buttons for forms
        update_btn = st.form_submit_button("✅ Update Job")
        delete_btn = st.form_submit_button("🗑 Delete Job")

        # ================= UPDATE =================
        if update_btn:

            cursor.execute(
                """
                UPDATE jobs
                SET job_date=?,
                    staff_name=?,
                    branch=?,
                    district_unit=?,
                    task=?,
                    reset_for=?,
                    description=?,
                    status=?
                WHERE id=?
                """,
                (
                    str(job_date),
                    staff_name,
                    branch,
                    district_unit,
                    task,
                    reset_for if task == "User account / AD reset" else None,
                    description,
                    status,
                    selected_job_id,
                ),
            )

            conn.commit()
            st.success("Job updated successfully ✅")


        # ================= DELETE =================
        if delete_btn:

            if st.session_state.role == "admin":

                cursor.execute(
                    "DELETE FROM jobs WHERE id=?",
                    (selected_job_id,),
                )
                conn.commit()
                st.success("Job deleted successfully 🗑")


            else:
                st.error("Only admin can delete jobs.")
def add_content_with_auto_slide(prs, title, lines, layout_index=1, max_lines=12):
    """
    Automatically creates new slides if content exceeds max_lines.
    """

    slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
    slide.shapes.title.text = title

    text_frame = slide.placeholders[1].text_frame
    text_frame.clear()

    line_count = 0

    for line in lines:

        if line_count >= max_lines:
            # Create new slide and continue
            slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
            slide.shapes.title.text = f"{title} (Cont.)"
            text_frame = slide.placeholders[1].text_frame
            text_frame.clear()
            line_count = 0

        p = text_frame.add_paragraph()
        p.text = line
        p.level = 0

        line_count += 1


def monthly_report_page():
    import calendar
    import datetime
    import os
    from pptx import Presentation

    st.markdown("## 📅 Monthly Report Compiler")

    df = pd.read_sql("SELECT * FROM jobs", conn)

    if df.empty:
        st.info("No data available.")
        return

    df["job_date"] = pd.to_datetime(df["job_date"])

    # ================= MONTH SELECTOR =================
    current_year = datetime.date.today().year
    current_month = datetime.date.today().month

    month_names = list(calendar.month_name)[1:]

    selected_month_name = st.selectbox(
        "Select Month",
        month_names,
        index=current_month - 1
    )

    selected_month = month_names.index(selected_month_name) + 1

    selected_year = st.selectbox(
        "Select Year",
        [current_year - 2, current_year - 1, current_year],
        index=2
    )

    df_month = df[
        (df["job_date"].dt.month == selected_month) &
        (df["job_date"].dt.year == selected_year)
    ]

    if df_month.empty:
        st.warning("No records for selected month.")
        return

    # ================= SUMMARY ON PAGE =================

    total_jobs = len(df_month)
    completed_jobs = len(df_month[df_month["status"] == "Completed"])
    ad_resets = len(df_month[df_month["task"] == "User account / AD reset"])

    st.markdown("### 📊 Monthly Summary")

    col1, col2, col3 = st.columns(3)
    col1.metric("Total Jobs", total_jobs)
    col2.metric("Completed", completed_jobs)
    col3.metric("AD Resets", ad_resets)


    # ================= TASK CATEGORIES =================

    task_categories = {
        "Active Directory (AD) Activities": [
            "User account / AD reset",
            "User account creation",
            "User account unlock",
            "Password reset"
        ],
        "Network & Connectivity Issues": [
            "Network troubleshooting",
            "Internet issue",
            "LAN issue",
            "IP configuration"
        ],
        "Hardware Support": [
            "PC repair",
            "Printer issue",
            "Hardware replacement"
        ],
        "Software Support": [
            "Software installation",
            "Application error",
            "System configuration"
        ]
    }

    # ================= GENERATE PPT =================

    if st.button("📥 Generate Monthly PPT"):

        template_path = "ion_boardroom_template.pptx"

        if not os.path.exists(template_path):
            st.error("monthly_template.pptx not found in project folder.")
            return

        prs = Presentation(template_path)

        # ========== COVER SLIDE ==========
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Monthly IS Activity Report"

        subtitle = slide.placeholders[1]
        subtitle.text = f"{selected_month_name} {selected_year}"

        # ========== SUMMARY SLIDE ==========
        summary_lines = [
            f"Total Jobs: {total_jobs}",
            f"Completed Jobs: {completed_jobs}",
            f"AD Resets: {ad_resets}"
        ]

        add_content_with_auto_slide(
            prs,
            "Monthly Summary",
            summary_lines,
            max_lines=10
        )

        # ========== TASK BREAKDOWN ==========
        task_counts = df_month["task"].value_counts()
        task_lines = [
            f"{task} : {count}"
            for task, count in task_counts.items()
        ]

        add_content_with_auto_slide(
            prs,
            "Task Breakdown",
            task_lines,
            max_lines=12
        )

        # ========== BRANCH ACTIVITY ==========
        branch_counts = df_month["branch"].value_counts()
        branch_lines = [
            f"{branch} : {count}"
            for branch, count in branch_counts.items()
        ]

        add_content_with_auto_slide(
            prs,
            "Branch Activity",
            branch_lines,
            max_lines=12
        )

        # ========== CATEGORIZED DETAILED ACTIVITIES ==========

        categorized_lines = []

        for category, task_list in task_categories.items():

            category_data = df_month[df_month["task"].isin(task_list)]

            if not category_data.empty:

                # Add category header
                categorized_lines.append(f"=== {category} ===")

                for _, row in category_data.iterrows():

                    reset_text = ""
                    if "reset_for" in df_month.columns and row.get("reset_for"):
                        reset_text = f" (Reset For: {row['reset_for']})"

                    description = row.get("description", "")

                    line = (
                        f"- {row['task']}{reset_text} "
                        f"({row['branch']}) "
                        f"[{description}]"
                    )

                    categorized_lines.append(line)

                categorized_lines.append("")


        other_tasks = df_month[
            ~df_month["task"].isin(
                [task for tasks in task_categories.values() for task in tasks]
            )
        ]

        if not other_tasks.empty:

            categorized_lines.append("=== Other Activities ===")

            for _, row in other_tasks.iterrows():
                description = row.get("description", "")

                line = (
                    f"- {row['task']} "
                    f"({row['branch']}) "
                    f"[{description}]"
                )

                categorized_lines.append(line)

        add_content_with_auto_slide(
            prs,
            "Detailed Activities",
            categorized_lines,
            max_lines=8
        )

        # ========== SAVE FILE ==========
        file_name = f"Monthly_Report_{selected_month}_{selected_year}.pptx"
        prs.save(file_name)

        with open(file_name, "rb") as f:
            st.download_button(
                "Download Monthly PPT",
                f,
                file_name=file_name
            )

def weekly_report_page():

    st.header("📊 Weekly Report")

    today = date.today()
    start_week = today - timedelta(days=today.weekday())
    end_week = start_week + timedelta(days=6)

    start_date = st.date_input("Week Start", start_week)
    end_date = st.date_input("Week End", end_week)

    # ------------------ LOAD DATA ------------------
    df_all = pd.read_sql(
        "SELECT * FROM jobs WHERE job_date BETWEEN ? AND ?",
        conn,
        params=(start_date, end_date),
    )

    # Staff can only see their jobs
    if st.session_state.role == "staff":
        df_view = df_all[
            df_all["staff_name"].apply(
                lambda x: x.replace(" ", "").lower() == st.session_state.username
            )
        ]
    else:
        df_view = df_all

    # ------------------ DISPLAY DATA ------------------
    if df_view.empty:
        st.info("No data for selected week.")
    else:
        st.subheader("📋 Weekly Data")
        st.dataframe(
            df_view[[
                "id",
                "job_date",
                "staff_name",
                "branch",
                "district_unit",
                "category",
                "task",
                "reset_for",
                "description",
                "status"
            ]]
        )

        # ------------------ STATUS SUMMARY ------------------
        st.subheader("📌 Status Summary")
        completed_count = len(df_view[df_view["status"] == "Completed"])
        inprogress_count = len(df_view[df_view["status"] == "In Progress"])

        col1, col2 = st.columns(2)
        col1.metric("Completed Jobs", completed_count)
        col2.metric("In Progress Jobs", inprogress_count)

        # ------------------ TASK SUMMARY ------------------
        st.subheader("📊 Task Type Comparison")
        task_summary = df_view["task"].value_counts().reset_index()
        task_summary.columns = ["Task Type", "Total Jobs"]
        st.dataframe(task_summary)

        # ------------------ CATEGORY SUMMARY ------------------
        st.subheader("📊 Category Comparison")
        category_summary = df_view["category"].value_counts().reset_index()
        category_summary.columns = ["Category", "Total Jobs"]
        st.dataframe(category_summary)

    # ==========================================================
    # =================== PPT EXPORT ============================
    # ==========================================================

    st.markdown("---")
    st.subheader("📤 Compile Weekly Report (PPT)")

    compiled_by = st.text_input(
        "Compiled By",
        value=st.session_state.username.capitalize()
    )

    if st.button("📤 Compile & Export PPT"):

        if df_all.empty:
            st.warning("No jobs to compile for this week.")
        else:
            prs = Presentation("ion_boardroom_template.pptx")

            # ---------- COVER SLIDE ----------
            cover_slide = prs.slides[0]
            cover_slide.shapes.title.text = ORGANIZATION
            cover_slide.placeholders[1].text = (
                f"{DISTRICT}\nIS Weekly Report\n\n"
                f"{start_date} to {end_date}\n\n"
                f"Compiled By: {compiled_by}"
            )

            layout = prs.slide_layouts[1]

            # ---------- STATUS SUMMARY SLIDE ----------
            summary_slide = prs.slides.add_slide(layout)
            summary_slide.shapes.title.text = "Weekly Status Summary"
            tf = summary_slide.placeholders[1].text_frame
            tf.clear()

            p1 = tf.add_paragraph()
            p1.text = f"Total Completed Jobs: {len(df_all[df_all['status'] == 'Completed'])}"

            p2 = tf.add_paragraph()
            p2.text = f"Total In Progress Jobs: {len(df_all[df_all['status'] == 'In Progress'])}"

            # ---------- TASK SUMMARY SLIDE ----------
            task_slide = prs.slides.add_slide(layout)
            task_slide.shapes.title.text = "Task Type Comparison"
            tf_task = task_slide.placeholders[1].text_frame
            tf_task.clear()

            task_summary_all = df_all["task"].value_counts().reset_index()
            task_summary_all.columns = ["Task Type", "Total Jobs"]

            for _, row in task_summary_all.iterrows():
                p = tf_task.add_paragraph()
                p.text = f"{row['Task Type']} : {row['Total Jobs']} Jobs"

            # ---------- CATEGORY DETAIL SLIDES ----------
            for category in df_all["category"].unique():
                slide_cat = prs.slides.add_slide(layout)
                slide_cat.shapes.title.text = category
                tf_cat = slide_cat.placeholders[1].text_frame
                tf_cat.clear()

                category_df = df_all[df_all["category"] == category]

                for _, row in category_df.iterrows():
                    district_unit_text = (
                        f" ({row['district_unit']})"
                        if row["district_unit"] else ""
                    )

                    p = tf_cat.add_paragraph()
                    p.text = (
                        f"- {row['task']}  "
                        f"({row['branch']}){district_unit_text} "
                        f"[{row['description']}]"
                    )

            filename = f"IS_Weekly_Report_{end_date}.pptx"
            prs.save(filename)

            with open(filename, "rb") as f:
                st.download_button(
                    "⬇ Download PPT",
                    f,
                    file_name=filename
                )

    # ==========================================================
    # =================== AD EXCEL REPORT ======================
    # ==========================================================

    st.markdown("---")
    st.subheader("📥 Export AD Report (Excel Only)")

    default_compiler = get_staff_name(st.session_state.username)

    compiled_by_ad = st.text_input(
        "Compiled By (AD Report)",
        value=default_compiler,
        key="compiled_by_ad"
    )

    if st.button("📥 Generate AD Report Excel"):

        df_ad = df_all[df_all["task"] == "User account / AD reset"]

        if df_ad.empty:
            st.warning("No AD reset jobs found for this week.")
        else:
            from openpyxl import Workbook
            from openpyxl.worksheet.table import Table, TableStyleInfo
            from openpyxl.utils import get_column_letter

            filename = f"AD_Report_{end_date}.xlsx"

            # Create workbook & sheet
            wb = Workbook()
            ws = wb.active
            ws.title = "AD Weekly Report"

            # Prepare export data
            df_ad_export = df_ad[[
                "job_date",
                "staff_name",
                "branch",
                "district_unit",
                "reset_for",
                "description",
                "status"
            ]].copy()

            df_ad_export["Compiled By"] = compiled_by_ad
            df_ad_export["Report Period"] = f"{start_date} to {end_date}"

            # Write headers
            ws.append(list(df_ad_export.columns))

            # Write data
            for row in df_ad_export.itertuples(index=False):
                ws.append(row)

            # Create Excel table
            last_row = ws.max_row
            last_col = ws.max_column
            table_range = f"A1:{get_column_letter(last_col)}{last_row}"

            table = Table(displayName="ADReportTable", ref=table_range)

            style = TableStyleInfo(
                name="TableStyleMedium9",
                showFirstColumn=False,
                showLastColumn=False,
                showRowStripes=True,
                showColumnStripes=False,
            )

            table.tableStyleInfo = style
            ws.add_table(table)

            # Auto-fit column width
            for col in ws.columns:
                max_length = 0
                col_letter = get_column_letter(col[0].column)
                for cell in col:
                    try:
                        if cell.value:
                            max_length = max(max_length, len(str(cell.value)))
                    except:
                        pass
                adjusted_width = max_length + 2
                ws.column_dimensions[col_letter].width = adjusted_width

            wb.save(filename)

            with open(filename, "rb") as f:
                st.download_button(
                    label="⬇ Download AD Excel Report",
                    data=f,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                )
def password_manager_page():
    st.markdown("### 🔑 Password Manager")

    if st.session_state.role == "staff":
        selected_user = st.session_state.username
        st.text(f"Your username: {selected_user}")
    else:
        users = [u[0] for u in cursor.execute("SELECT username FROM users").fetchall()]
        selected_user = st.selectbox("Select User", users)

    with st.form(f"pw_form_{selected_user}"):
        new_password = st.text_input("New Password", type="password")
        confirm_password = st.text_input("Confirm Password", type="password")
        submit_btn = st.form_submit_button("Change Password")

        if submit_btn:
            if new_password != confirm_password:
                st.error("Passwords do not match")
            else:
                if st.session_state.role == "staff" and selected_user != st.session_state.username:
                    st.error("You cannot change other users' passwords")
                else:
                    cursor.execute("UPDATE users SET password=? WHERE username=?", (new_password, selected_user))
                    conn.commit()
                    st.success("Password updated successfully")


# ================= TOP NAV =================
def top_nav():
    pages = ["🏠 Dashboard","📝 Job Entry","✏️ Update/Delete","📊 Weekly Report","🔑 Password Manager"," Monthly Report"]
    cols = st.columns(len(pages))
    for i, page in enumerate(pages):
        if cols[i].button(page, key=f"nav_{page}"):
            st.session_state.page = page

# ================= PAGE ROUTER =================
if not st.session_state.logged_in:
    login_page()
else:
    st.set_page_config(layout="wide")
    top_nav()
    if st.session_state.page=="🏠 Dashboard":
        dashboard_page()
    elif st.session_state.page=="📝 Job Entry":
        job_entry_page()
    elif st.session_state.page=="✏️ Update/Delete":
        update_delete_page()
    elif st.session_state.page=="📊 Weekly Report":
        weekly_report_page()
    elif st.session_state.page==" Monthly Report":
        monthly_report_page()
    elif st.session_state.page=="🔑 Password Manager":
        password_manager_page()